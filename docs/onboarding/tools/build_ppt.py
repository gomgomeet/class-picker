# -*- coding: utf-8 -*-
"""클로드 코드 온보딩 PPT 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image
import os

ASSETS = '/home/user/class-picker/docs/onboarding/assets'
SHOTS = '/home/user/class-picker/docs/onboarding/assets/screenshots'
OUT = '/home/user/class-picker/docs/onboarding/claude-code-onboarding.pptx'

BG    = RGBColor(0x17, 0x17, 0x1B)
CARD  = RGBColor(0x22, 0x22, 0x29)
CORAL = RGBColor(0xD9, 0x77, 0x57)
TXT   = RGBColor(0xEC, 0xEC, 0xEA)
DIM   = RGBColor(0x9B, 0x9B, 0x93)
GREEN = RGBColor(0x8F, 0xBC, 0x6F)
YELL  = RGBColor(0xD9, 0xB4, 0x5B)

SW, SH = Inches(13.333), Inches(7.5)
FONT = 'Malgun Gothic'
MONO = 'Consolas'

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]

def _ea(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)

def add_run(p, text, size, color=TXT, bold=False, font=FONT, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    _ea(r, FONT if font == FONT else font)
    return r

def tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf

def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def footer(s, num, label=''):
    tf = tb(s, Inches(0.45), Inches(7.05), Inches(9), Inches(0.35))
    p = tf.paragraphs[0]
    add_run(p, label, 11, DIM)
    tf2 = tb(s, Inches(12.3), Inches(7.05), Inches(0.7), Inches(0.35))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    add_run(p2, str(num), 11, DIM)

def accent_bar(s, x=Inches(0.5), y=Inches(0.62), w=Inches(0.55), h=Inches(0.07)):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = CORAL
    sh.line.fill.background()
    return sh

def title(s, text, sub=None):
    accent_bar(s)
    tf = tb(s, Inches(0.45), Inches(0.78), Inches(12.4), Inches(0.9))
    p = tf.paragraphs[0]
    add_run(p, text, 30, TXT, bold=True)
    if sub:
        tf2 = tb(s, Inches(0.47), Inches(1.42), Inches(12.4), Inches(0.5))
        add_run(tf2.paragraphs[0], sub, 15, DIM)

def bullets(s, items, x, y, w, size=17, gap=8):
    tf = tb(s, x, y, w, Inches(5))
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        if isinstance(it, tuple) and it and it[0] == 'sub':
            add_run(p, '     ' + it[1], size - 3, DIM)
            continue
        segs = it if isinstance(it, list) else [(it, TXT, False)]
        add_run(p, '•  ', size, CORAL, bold=True)
        for seg in segs:
            t, c, b = (seg + (TXT, False))[:3] if isinstance(seg, tuple) else (seg, TXT, False)
            add_run(p, t, size, c, bold=b)
    return tf

def pic(s, name, x=None, y=None, w=None, h=None, base=None):
    path = os.path.join(base or ASSETS, f'{name}.png')
    iw, ih = Image.open(path).size
    if w and not h:
        h = w * ih / iw
    if h and not w:
        w = h * iw / ih
    if x is None: x = (SW - w) / 2
    if y is None: y = Inches(1.9)
    return s.shapes.add_picture(path, int(x), int(y), int(w), int(h))

def card(s, x, y, w, h, border=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = CARD
    if border:
        sh.line.color.rgb = border; sh.line.width = Pt(1.2)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

N = 0
def num():
    global N
    N += 1
    return N

# ══════════════ 1. 표지 ══════════════
s = new_slide(); num()
tf = tb(s, Inches(0.7), Inches(0.85), Inches(12), Inches(2.2))
p = tf.paragraphs[0]
add_run(p, '클로드 코드 시작하기', 44, TXT, bold=True)
p2 = tf.add_paragraph(); p2.space_before = Pt(10)
add_run(p2, '터미널에서 만나는 AI 코딩 에이전트 — 처음 사용하는 분들을 위한 온보딩', 18, DIM)
pic(s, 'cover-terminal', x=Inches(2.47), y=Inches(2.55), w=Inches(8.4))
tf = tb(s, Inches(0.7), Inches(7.0), Inches(12), Inches(0.4))
add_run(tf.paragraphs[0], 'Claude Code Onboarding  ·  v1.0  ·  2026', 12, DIM)

# ══════════════ 2. 로드맵 ══════════════
s = new_slide()
title(s, '오늘 배울 것 — 3단계 로드맵', 'Level 1만 따라 해도 오늘부터 쓸 수 있습니다. Level 2·3은 필요할 때 돌아오세요.')
pic(s, 'diagram-levels-roadmap', y=Inches(1.95), w=Inches(9.4))
footer(s, num(), '인트로')

# ══════════════ 3. 클로드 코드란? ══════════════
s = new_slide()
title(s, '클로드 코드란?')
card(s, Inches(0.5), Inches(1.75), Inches(12.33), Inches(1.25), border=CORAL)
tf = tb(s, Inches(0.9), Inches(1.95), Inches(11.6), Inches(0.9), MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]
add_run(p, '“터미널에서 실행되는 Anthropic의 ', 20, TXT)
add_run(p, 'AI 코딩 에이전트', 20, CORAL, bold=True)
add_run(p, ' — 자연어로 말하면 코드베이스를 이해하고 직접 일합니다.”', 20, TXT)
bullets(s, [
    [('말로 요청 ', TXT, True), ('— "다크 모드 추가해줘" 처럼 목표만 말하면 됩니다', TXT, False)],
    [('직접 작업 ', TXT, True), ('— 파일을 읽고, 수정하고, 명령을 실행하고, 테스트합니다', TXT, False)],
    [('Git까지 ', TXT, True), ('— 커밋 메시지 작성, 브랜치, PR 생성까지 한 흐름으로', TXT, False)],
    [('승인 기반 ', TXT, True), ('— 모든 변경은 내가 확인하고 허락해야 실행됩니다', TXT, False)],
], Inches(0.7), Inches(3.35), Inches(12), size=18, gap=14)
footer(s, num(), '인트로')

# ══════════════ 4. 무엇이 다른가 ══════════════
s = new_slide()
title(s, '자동완성과 무엇이 다른가요?')
pic(s, 'diagram-agent-vs-autocomplete', y=Inches(1.8), w=Inches(10.6))
footer(s, num(), '인트로')

def section(sec_label, big, desc, goals):
    s = new_slide()
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.0), Inches(0.18), Inches(3.5))
    sh.fill.solid(); sh.fill.fore_color.rgb = CORAL; sh.line.fill.background()
    tf = tb(s, Inches(0.85), Inches(2.05), Inches(11.5), Inches(1.0))
    add_run(tf.paragraphs[0], sec_label, 17, CORAL, bold=True)
    tf = tb(s, Inches(0.85), Inches(2.45), Inches(11.5), Inches(1.3))
    add_run(tf.paragraphs[0], big, 40, TXT, bold=True)
    tf = tb(s, Inches(0.87), Inches(3.5), Inches(11.5), Inches(0.6))
    add_run(tf.paragraphs[0], desc, 17, DIM)
    bullets(s, goals, Inches(0.87), Inches(4.25), Inches(11.5), size=16, gap=8)
    footer(s, num(), sec_label)
    return s

# ══════════════ 5. 섹션: 설치 ══════════════
section('설치 & 첫 실행', '5분이면 준비 끝', '설치하고, 로그인하고, 첫 인사를 나눠 봅시다.',
        ['설치 명령 한 줄과 로그인', '첫 대화 시작하는 법', '권한 승인 흐름 이해하기'])

# ══════════════ 6. 설치하기 ══════════════
s = new_slide()
title(s, '설치하기')
bullets(s, [
    [('요구사항: ', TXT, True), ('macOS · Linux · Windows, Node.js 18+', TXT, False)],
    [('npm 설치:', TXT, True)],
    ('sub', 'npm install -g @anthropic-ai/claude-code'),
    [('프로젝트 폴더에서 실행:', TXT, True)],
    ('sub', 'cd 내-프로젝트  →  claude'),
    [('업데이트는 자동 ', TXT, True), ('— 항상 최신 기능을 사용합니다', DIM, False)],
], Inches(0.6), Inches(2.0), Inches(5.6), size=16, gap=11)
pic(s, 'term-install', x=Inches(6.45), y=Inches(1.85), w=Inches(6.4))
footer(s, num(), '설치 & 첫 실행')

# ══════════════ 7. 로그인 ══════════════
s = new_slide()
title(s, '로그인과 요금')
bullets(s, [
    [('첫 실행 시 브라우저가 열리며 ', TXT, False), ('Claude 계정으로 로그인', TXT, True), ('합니다', TXT, False)],
    [('두 가지 이용 방법', TXT, True)],
    ('sub', '① Claude Pro / Max 구독 — 구독에 포함, 추가 비용 없음'),
    ('sub', '② Claude API 키 — 쓴 만큼 과금 (팀·자동화에 적합)'),
    [('/status', CORAL, True), (' 로 로그인 상태·현재 모델을 언제든 확인', TXT, False)],
    [('/model', CORAL, True), (' 로 모델 변경 — 복잡한 작업은 Opus, 빠른 작업은 Sonnet', TXT, False)],
], Inches(0.6), Inches(2.0), Inches(12), size=17, gap=13)
footer(s, num(), '설치 & 첫 실행')

# ══════════════ 8. 첫 대화 ══════════════
s = new_slide()
title(s, '첫 대화 — 이렇게 물어보세요', '파일명이나 함수명을 몰라도 됩니다. Claude가 스스로 찾아냅니다.')
ex = [('“이 프로젝트 구조를 설명해줘”', '처음 보는 코드베이스 파악'),
      ('“빌드하고 실행하려면 어떻게 해?”', '문서 없는 프로젝트도 OK'),
      ('“로그인 처리하는 코드가 어디 있어?”', '키워드만으로 코드 탐색')]
x = Inches(0.6)
for q, d in ex:
    card(s, x, Inches(2.1), Inches(3.95), Inches(2.1))
    tf = tb(s, x + Inches(0.28), Inches(2.35), Inches(3.4), Inches(1.2))
    add_run(tf.paragraphs[0], q, 16, TXT, bold=True)
    tf = tb(s, x + Inches(0.28), Inches(3.55), Inches(3.4), Inches(0.5))
    add_run(tf.paragraphs[0], d, 13, DIM)
    x += Inches(4.15)
tf = tb(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.6))
p = tf.paragraphs[0]
add_run(p, 'TIP  ', 14, GREEN, bold=True)
add_run(p, '읽기만 하는 질문은 코드를 바꾸지 않습니다 — 부담 없이 많이 물어보세요.', 15, DIM)
footer(s, num(), '설치 & 첫 실행')

# ══════════════ 9. 권한 시스템 ══════════════
s = new_slide()
title(s, '권한 시스템 — 허락 없이는 아무것도 바꾸지 않습니다')
bullets(s, [
    ['파일 수정·명령 실행 전 ', ('항상 먼저 물어봅니다', TXT, True)],
    [('diff(변경 내용)를 눈으로 확인', TXT, True), (' 후 승인 / 거부', TXT, False)],
    ['반복 승인이 번거로우면 ', ('“allow all”', CORAL, True), ('로 세션 동안 자동 승인', TXT, False)],
    [('Esc', CORAL, True), (' 로 언제든 중단 · ', TXT, False), ('Esc Esc', CORAL, True), (' 로 이전 시점 복원', TXT, False)],
], Inches(0.6), Inches(2.1), Inches(5.7), size=16, gap=13)
pic(s, 'term-diff-approval', x=Inches(6.45), y=Inches(1.8), w=Inches(6.4))
footer(s, num(), '설치 & 첫 실행')

# ══════════════ 10. 섹션: Level 1 ══════════════
section('Level 1', '시작하기', '첫 30분 — 질문하고, 수정하고, 기본 명령어에 익숙해집니다.',
        ['코드베이스에 자유롭게 질문하기', '자연어로 파일 수정 요청하기', '필수 슬래시 커맨드 5개'])

# ══════════════ 11. 코드베이스 질문 ══════════════
s = new_slide()
title(s, '코드베이스 질문하기')
bullets(s, [
    ['Claude가 필요한 파일을 ', ('스스로 찾아 읽고', TXT, True), ' 답합니다'],
    ['새 팀 합류, 오래된 코드, 남의 코드 파악에 최고',],
    [('활용 예', TXT, True)],
    ('sub', '“이 함수는 어디서 호출돼?”'),
    ('sub', '“결제 로직의 흐름을 단계별로 설명해줘”'),
    ('sub', '“이 에러가 왜 나는지 찾아줘”'),
], Inches(0.6), Inches(2.1), Inches(5.7), size=16, gap=11)
pic(s, 'term-codebase-question', x=Inches(6.45), y=Inches(1.8), w=Inches(6.4))
footer(s, num(), 'Level 1 · 시작하기')

# ══════════════ 12. 파일 수정 ══════════════
s = new_slide()
title(s, '파일 수정 요청하기', '원하는 "결과"를 말하세요 — 방법은 Claude가 찾습니다.')
rows = [('아쉬운 요청', '좋은 요청', DIM, CORAL),
        ('“코드 좀 고쳐줘”', '“pick()이 같은 학생을 두 번 안 뽑게 고쳐줘”', None, None),
        ('“버그 있어”', '“뽑기 버튼을 두 번 누르면 결과가 사라져. 원인 찾아서 고쳐줘”', None, None),
        ('“예쁘게 해줘”', '“버튼을 파란색 계열로, 모서리 둥글게 바꿔줘”', None, None)]
y = Inches(2.0)
for i, (a, b, ca, cb) in enumerate(rows):
    hdr = i == 0
    card(s, Inches(0.6), y, Inches(5.5), Inches(0.85), border=None)
    card(s, Inches(6.3), y, Inches(6.4), Inches(0.85), border=CORAL if hdr else None)
    tf = tb(s, Inches(0.85), y, Inches(5.1), Inches(0.85), MSO_ANCHOR.MIDDLE)
    add_run(tf.paragraphs[0], a, 15 if not hdr else 15, ca or (DIM if not hdr else DIM), bold=hdr)
    tf = tb(s, Inches(6.55), y, Inches(6.0), Inches(0.85), MSO_ANCHOR.MIDDLE)
    add_run(tf.paragraphs[0], b, 15, cb or TXT, bold=hdr)
    y += Inches(1.02)
tf = tb(s, Inches(0.6), y + Inches(0.1), Inches(12), Inches(0.5))
p = tf.paragraphs[0]
add_run(p, 'TIP  ', 14, GREEN, bold=True)
add_run(p, '대상 · 원하는 동작 · 제약 조건, 세 가지가 들어가면 한 번에 됩니다.', 15, DIM)
footer(s, num(), 'Level 1 · 시작하기')

# ══════════════ 13. 슬래시 커맨드 ══════════════
s = new_slide()
title(s, '필수 슬래시 커맨드', '입력창에 / 를 치면 전체 목록이 나옵니다.')
bullets(s, [
    [('/help', CORAL, True), ('  사용법과 명령어 목록', TXT, False)],
    [('/init', CORAL, True), ('  CLAUDE.md 자동 생성 (Level 2에서 자세히)', TXT, False)],
    [('/clear', CORAL, True), ('  새 작업 시작 전 대화 초기화', TXT, False)],
    [('/compact', CORAL, True), ('  긴 대화 요약으로 컨텍스트 절약', TXT, False)],
    [('/model', CORAL, True), ('  Claude 모델 변경', TXT, False)],
], Inches(0.6), Inches(2.1), Inches(5.7), size=16, gap=12)
pic(s, 'term-slash-commands', x=Inches(6.45), y=Inches(1.85), w=Inches(6.4))
footer(s, num(), 'Level 1 · 시작하기')

# ══════════════ 14. 대화 제어 ══════════════
s = new_slide()
title(s, '대화 제어하기')
items = [('Esc', '진행 중인 작업 즉시 중단', '방향이 이상하다 싶으면 바로 멈추세요'),
         ('Esc Esc', '이전 시점으로 되돌리기', '대화·코드를 체크포인트로 복원 (rewind)'),
         ('claude -c', '마지막 대화 이어가기', '터미널을 껐다 켜도 이어서 작업'),
         ('claude -r', '과거 세션 목록에서 선택해 재개', '어제 하던 작업 다시 열기')]
y = Inches(2.0)
for k, t, d in items:
    card(s, Inches(0.6), y, Inches(12.1), Inches(1.0))
    tf = tb(s, Inches(0.95), y, Inches(2.6), Inches(1.0), MSO_ANCHOR.MIDDLE)
    add_run(tf.paragraphs[0], k, 17, CORAL, bold=True, font=MONO)
    tf = tb(s, Inches(3.7), y, Inches(4.3), Inches(1.0), MSO_ANCHOR.MIDDLE)
    add_run(tf.paragraphs[0], t, 16, TXT, bold=True)
    tf = tb(s, Inches(8.1), y, Inches(4.4), Inches(1.0), MSO_ANCHOR.MIDDLE)
    add_run(tf.paragraphs[0], d, 13, DIM)
    y += Inches(1.16)
footer(s, num(), 'Level 1 · 시작하기')

# ══════════════ 15. 라이브 데모 1 ══════════════
s = new_slide()
title(s, '라이브 데모 ①  —  class-picker에 기능 추가하기')
bullets(s, [
    [('시나리오', TXT, True)],
    ('sub', '1) “이 프로젝트 구조 설명해줘” — 코드베이스 파악'),
    ('sub', '2) “뽑기 결과가 화면에 크게 표시되게 해줘” — 수정 승인 흐름'),
    ('sub', '3) 브라우저 새로고침으로 결과 확인'),
    ['예상 소요: ', ('5분', CORAL, True)],
], Inches(0.55), Inches(2.0), Inches(4.9), size=16, gap=10)
# 데모 재현 애니메이션 GIF (발표 시 자동 재생 · 라이브 데모 백업)
gif1 = s.shapes.add_picture(os.path.join(ASSETS, 'demo1-backup.gif'),
                            Inches(5.55), Inches(1.9), width=Inches(7.3))
tf = tb(s, Inches(5.55), Inches(6.6), Inches(7.3), Inches(0.4))
add_run(tf.paragraphs[0], '▶ 데모 재현 애니메이션 (발표 화면에서 자동 재생) · assets/demo1-backup.gif', 11, CORAL)
# 실제 결과 화면 인셋
pic(s, 'shot-19-170006-app-winner', x=Inches(0.6), y=Inches(4.5), h=Inches(2.25), base=SHOTS)
tf = tb(s, Inches(3.05), Inches(4.8), Inches(2.4), Inches(1.5), MSO_ANCHOR.MIDDLE)
add_run(tf.paragraphs[0], '← 실제 수업에서\n   만든 결과 화면', 12, DIM)
footer(s, num(), 'Level 1 · 시작하기')

# ══════════════ 16. Level 1 정리 ══════════════
s = new_slide()
title(s, 'Level 1 정리 — 오늘부터 이렇게 쓰세요')
bullets(s, [
    ['모르는 코드는 ', ('일단 질문', CORAL, True), ' — 읽기 질문은 코드를 바꾸지 않습니다'],
    ['수정 요청은 ', ('결과 중심으로 구체적으로', CORAL, True)],
    ['diff를 확인하는 습관 — ', ('승인 전에 꼭 읽기', TXT, True)],
    [('실습 과제', GREEN, True), ('  내 프로젝트에서 ① 구조 질문 ② 작은 수정 1건 ③ Esc Esc로 되돌려 보기', TXT, False)],
], Inches(0.7), Inches(2.3), Inches(12), size=18, gap=18)
footer(s, num(), 'Level 1 · 시작하기')

# ══════════════ 17. 섹션: Level 2 ══════════════
section('Level 2', '업무에 붙이기', '반복 설명을 없애고, 안전하게 맡기고, Git까지 연결합니다.',
        ['CLAUDE.md로 프로젝트 기억시키기', 'Plan Mode로 큰 작업 안전하게', '커밋 · PR까지 한 흐름으로'])

# ══════════════ 18. CLAUDE.md ══════════════
s = new_slide()
title(s, 'CLAUDE.md — 프로젝트를 기억시키는 파일', '/init 한 번이면 초안이 자동 생성됩니다. 매 세션 자동으로 로드됩니다.')
pic(s, 'diagram-claude-md', y=Inches(1.85), w=Inches(9.9))
footer(s, num(), 'Level 2 · 업무에 붙이기')

# ══════════════ 19. Plan Mode ══════════════
s = new_slide()
title(s, 'Plan Mode — 실행 전에 계획부터')
bullets(s, [
    [('Shift+Tab', CORAL, True), (' 으로 전환 — 계획을 먼저 보고 승인', TXT, False)],
    ['이럴 때 쓰세요:'],
    ('sub', '· 여러 파일을 건드리는 큰 변경'),
    ('sub', '· 리팩터링, 마이그레이션'),
    ('sub', '· 처음 맡겨보는 종류의 작업'),
    ['계획이 마음에 안 들면 ', ('그 자리에서 피드백', TXT, True), ' — 코드는 아직 안 바뀌었으니 안전'],
], Inches(0.6), Inches(2.05), Inches(5.7), size=16, gap=10)
pic(s, 'term-plan-mode', x=Inches(6.45), y=Inches(1.8), w=Inches(6.4))
footer(s, num(), 'Level 2 · 업무에 붙이기')

# ══════════════ 20. 권한 모드 ══════════════
s = new_slide()
title(s, '권한 모드 3가지', 'Shift+Tab 으로 순환합니다. 상황에 맞는 모드를 고르세요.')
modes = [('기본 모드', '수정·명령마다 승인', '처음 쓸 때 · 중요한 코드', TXT),
         ('Auto-accept', '파일 수정 자동 승인', '반복 작업 · 신뢰가 쌓인 뒤', YELL),
         ('Plan Mode', '계획만 세우고 실행 안 함', '큰 변경의 사전 검토', CORAL)]
x = Inches(0.6)
for name, what, when, c in modes:
    card(s, x, Inches(2.1), Inches(3.95), Inches(3.3), border=c if c != TXT else None)
    tf = tb(s, x + Inches(0.3), Inches(2.4), Inches(3.35), Inches(0.6))
    add_run(tf.paragraphs[0], name, 19, c, bold=True)
    tf = tb(s, x + Inches(0.3), Inches(3.15), Inches(3.35), Inches(0.9))
    add_run(tf.paragraphs[0], what, 15, TXT)
    tf = tb(s, x + Inches(0.3), Inches(4.35), Inches(3.35), Inches(0.9))
    p = tf.paragraphs[0]
    add_run(p, '추천 상황\n', 12, DIM, bold=True)
    add_run(p, when, 14, DIM)
    x += Inches(4.15)
tf = tb(s, Inches(0.6), Inches(5.75), Inches(12), Inches(0.6))
p = tf.paragraphs[0]
add_run(p, '주의  ', 14, YELL, bold=True)
add_run(p, '자동 승인 모드에서도 위험한 명령(삭제, 배포 등)은 다시 물어봅니다.', 15, DIM)
footer(s, num(), 'Level 2 · 업무에 붙이기')

# ══════════════ 21. Git ══════════════
s = new_slide()
title(s, 'Git 작업 맡기기')
pic(s, 'diagram-git-workflow', y=Inches(1.8), w=Inches(10.6))
footer(s, num(), 'Level 2 · 업무에 붙이기')

# ══════════════ 22. 이미지 활용 ══════════════
s = new_slide()
title(s, '이미지로 소통하기', '스크린샷을 터미널에 붙여넣거나(Ctrl+V) 파일을 드래그하세요.')
ex = [('디자인 시안 →', '“이 시안대로 페이지를 만들어줘”'),
      ('에러 스크린샷 →', '“이 화면처럼 깨져. 원인 찾아서 고쳐줘”'),
      ('완성 화면 비교 →', '“구현 결과가 시안과 다른 부분을 찾아줘”'),
      ('손그림 와이어프레임 →', '“이 스케치 레이아웃으로 잡아줘”')]
y = Inches(2.1)
for a, b in ex:
    card(s, Inches(0.6), y, Inches(12.1), Inches(0.95))
    tf = tb(s, Inches(0.95), y, Inches(3.6), Inches(0.95), MSO_ANCHOR.MIDDLE)
    add_run(tf.paragraphs[0], a, 16, CORAL, bold=True)
    tf = tb(s, Inches(4.7), y, Inches(7.8), Inches(0.95), MSO_ANCHOR.MIDDLE)
    add_run(tf.paragraphs[0], b, 16, TXT)
    y += Inches(1.1)
footer(s, num(), 'Level 2 · 업무에 붙이기')

# ══════════════ 22.5 배포하기 ══════════════
s = new_slide()
title(s, '만든 것을 세상에 공개하기 — 배포', '코드를 몰라도 됩니다. 무료이고, 학생들 폰에서도 열립니다.')
bullets(s, [
    [('① Netlify Drop', CORAL, True), ('  app.netlify.com/drop 에 폴더를 드래그 — 몇 초면 주소 완성', TXT, False)],
    [('② GitHub Pages', CORAL, True), ('  저장소 Settings → Pages — 오래 쓸 안정적인 주소', TXT, False)],
    ['둘 다 무료 · HTTPS — 주소만 알면 교실 어느 기기에서든 열림'],
    [('주의  ', YELL, True), ('공개 페이지에는 학생 실명을 코드에 넣지 마세요 — 수업 때 입력창에 직접 입력', TXT, False)],
    ['Claude에게 “배포하는 방법 알려줘”라고 물으면 단계별로 안내해 줍니다'],
], Inches(0.6), Inches(2.05), Inches(6.5), size=15, gap=12)
pic(s, 'shot-23-170453-deploy-guide', x=Inches(7.5), y=Inches(1.85), h=Inches(4.9), base=SHOTS)
tf = tb(s, Inches(7.5), Inches(6.8), Inches(5.3), Inches(0.4))
add_run(tf.paragraphs[0], '실제 수업에서 Claude가 안내한 배포 방법', 11, DIM)
footer(s, num(), 'Level 2 · 업무에 붙이기')

# ══════════════ 23. 컨텍스트 관리 ══════════════
s = new_slide()
title(s, '컨텍스트 관리 — 길어진 대화 다루기')
bullets(s, [
    ['대화가 길어지면 응답이 느려지고 초점이 흐려집니다'],
    [('/clear', CORAL, True), ('  주제가 바뀔 때 — 새 작업은 새 대화로', TXT, False)],
    [('/compact', CORAL, True), ('  이어가야 할 때 — 지금까지를 요약하고 계속', TXT, False)],
    [('어려운 문제엔 확장 사고 ', TXT, True), ('— 요청에 “think hard”를 붙이면 더 깊이 고민합니다', TXT, False)],
    [('작업 단위로 끊기', TXT, True), (' — 한 대화 = 한 가지 목표가 가장 정확합니다', TXT, False)],
], Inches(0.7), Inches(2.2), Inches(12), size=17, gap=15)
footer(s, num(), 'Level 2 · 업무에 붙이기')

# ══════════════ 24. 라이브 데모 2 ══════════════
s = new_slide()
title(s, '라이브 데모 ②  —  Plan Mode로 개선하고 PR까지')
bullets(s, [
    [('시나리오', TXT, True)],
    ('sub', '1) Shift+Tab → Plan Mode 전환'),
    ('sub', '2) “같은 학생이 연속으로 안 뽑히게 개선해줘” — 계획 검토'),
    ('sub', '3) 승인 → 구현 → “브랜치 만들어서 PR 올려줘”'),
    ('sub', '4) GitHub에서 PR 확인'),
    ['예상 소요: ', ('8분', CORAL, True)],
], Inches(0.55), Inches(2.0), Inches(4.9), size=16, gap=9)
# 데모 재현 애니메이션 GIF (발표 시 자동 재생 · 라이브 데모 백업)
s.shapes.add_picture(os.path.join(ASSETS, 'demo2-backup.gif'),
                     Inches(5.55), Inches(1.9), width=Inches(7.3))
tf = tb(s, Inches(5.55), Inches(6.6), Inches(7.3), Inches(0.4))
add_run(tf.paragraphs[0], '▶ 데모 재현 애니메이션 (발표 화면에서 자동 재생) · assets/demo2-backup.gif', 11, CORAL)
# 실제 앱 화면 인셋
pic(s, 'shot-14-165313-app-full-ui', x=Inches(0.6), y=Inches(4.55), h=Inches(2.2), base=SHOTS)
tf = tb(s, Inches(2.9), Inches(4.85), Inches(2.5), Inches(1.5), MSO_ANCHOR.MIDDLE)
add_run(tf.paragraphs[0], '← 실제 수업에서\n   만든 앱 화면', 12, DIM)
footer(s, num(), 'Level 2 · 업무에 붙이기')

# ══════════════ 25. Level 2 정리 ══════════════
s = new_slide()
title(s, 'Level 2 정리')
bullets(s, [
    [('/init', CORAL, True), (' 은 오늘 바로 — 팀 프로젝트라면 CLAUDE.md를 커밋해서 공유', TXT, False)],
    ['큰 작업은 ', ('Plan Mode가 기본값', CORAL, True), ' 이라고 생각하세요'],
    ['커밋 메시지 · PR 설명은 이제 Claude의 일'],
    [('실습 과제', GREEN, True), ('  ① /init 실행 후 CLAUDE.md 다듬기 ② Plan Mode로 작은 리팩터링 ③ PR 1건 생성', TXT, False)],
], Inches(0.7), Inches(2.3), Inches(12), size=18, gap=18)
footer(s, num(), 'Level 2 · 업무에 붙이기')

# ══════════════ 26. 섹션: Level 3 ══════════════
section('Level 3', '확장하기', '"이런 것도 된다" — 필요해졌을 때 돌아와서 깊게 보세요.',
        ['반복 작업 자동화 (커스텀 커맨드 · 스킬 · Hooks)', '서브에이전트로 병렬 작업', 'MCP로 외부 도구 연결 · 어디서든 사용'])

# ══════════════ 27. 커스텀 커맨드 & 스킬 ══════════════
s = new_slide()
title(s, '커스텀 슬래시 커맨드 & 스킬', '자주 하는 요청을 나만의 명령으로 저장하세요.')
bullets(s, [
    [('.claude/commands/deploy-check.md', CORAL, True), (' 파일을 만들면 → ', TXT, False), ('/deploy-check', CORAL, True), (' 명령 완성', TXT, False)],
    ['프롬프트를 파일로 저장하는 것 — 마크다운만 알면 됩니다'],
    [('스킬(Skills)', TXT, True), (' 은 더 큰 단위 — 절차·체크리스트·스크립트를 묶은 재사용 패키지', TXT, False)],
    [('활용 예', TXT, True)],
    ('sub', '/review-pr — 우리 팀 체크리스트로 코드 리뷰'),
    ('sub', '/release-note — 커밋 로그에서 릴리스 노트 생성'),
], Inches(0.7), Inches(2.1), Inches(12), size=17, gap=13)
footer(s, num(), 'Level 3 · 확장하기')

# ══════════════ 28. 서브에이전트 ══════════════
s = new_slide()
title(s, '서브에이전트 — 일을 나눠 맡기기')
bullets(s, [
    ['별도 컨텍스트에서 도는 ', ('보조 Claude', CORAL, True), ' — 역할·도구·모델을 지정해 정의'],
    ['본 대화를 어지럽히지 않고 큰 탐색·검증 작업을 위임'],
    [('활용 예', TXT, True)],
    ('sub', '코드 리뷰어 에이전트 — 변경마다 자동으로 버그 사냥'),
    ('sub', '테스트 러너 에이전트 — 실패한 테스트만 골라 수정 제안'),
    ('sub', '탐색 에이전트 — 넓은 코드베이스에서 관련 코드 수집'),
    [('.claude/agents/', CORAL, True), (' 폴더에 마크다운으로 정의합니다', TXT, False)],
], Inches(0.7), Inches(2.1), Inches(12), size=17, gap=12)
footer(s, num(), 'Level 3 · 확장하기')

# ══════════════ 29. MCP ══════════════
s = new_slide()
title(s, 'MCP — 외부 도구 연결')
pic(s, 'diagram-mcp', y=Inches(1.75), w=Inches(9.9))
footer(s, num(), 'Level 3 · 확장하기')

# ══════════════ 30. Hooks ══════════════
s = new_slide()
title(s, 'Hooks — 규칙을 자동으로 강제하기', 'Claude의 특정 동작 전후에 내 스크립트를 실행합니다.')
ex = [('파일 수정 후', '자동으로 포매터 실행', 'PostToolUse → prettier'),
      ('명령 실행 전', '위험한 명령 차단', 'PreToolUse → rm -rf 금지'),
      ('작업 완료 시', '알림 보내기', 'Stop → 슬랙/데스크톱 알림')]
x = Inches(0.6)
for a, b, c in ex:
    card(s, x, Inches(2.2), Inches(3.95), Inches(2.6))
    tf = tb(s, x + Inches(0.3), Inches(2.5), Inches(3.35), Inches(0.5))
    add_run(tf.paragraphs[0], a, 14, DIM, bold=True)
    tf = tb(s, x + Inches(0.3), Inches(3.0), Inches(3.35), Inches(0.9))
    add_run(tf.paragraphs[0], b, 17, TXT, bold=True)
    tf = tb(s, x + Inches(0.3), Inches(4.05), Inches(3.35), Inches(0.6))
    add_run(tf.paragraphs[0], c, 13, CORAL, font=MONO)
    x += Inches(4.15)
tf = tb(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.6))
p = tf.paragraphs[0]
add_run(p, '“부탁"이 아니라 "보장" — ', 15, TXT, bold=True)
add_run(p, 'CLAUDE.md의 규칙은 권고, Hooks는 강제입니다.', 15, DIM)
footer(s, num(), 'Level 3 · 확장하기')

# ══════════════ 31. 에코시스템 ══════════════
s = new_slide()
title(s, '어디서든 쓰기')
pic(s, 'diagram-ecosystem', y=Inches(1.8), w=Inches(10.2))
footer(s, num(), 'Level 3 · 확장하기')

# ══════════════ 32. Level 3 정리 ══════════════
s = new_slide()
title(s, 'Level 3 정리 — 언제 무엇을 꺼내 쓸까')
ex = [('같은 요청을 3번째 타이핑할 때', '커스텀 커맨드'),
      ('“매번 이것 좀 지켜줬으면” 할 때', 'CLAUDE.md → 안 지켜지면 Hooks'),
      ('작업이 크고 검증도 필요할 때', '서브에이전트'),
      ('Slack·Notion·DB가 필요한 작업', 'MCP'),
      ('자리를 떠나야 할 때', '웹(claude.ai/code) · GitHub Actions')]
y = Inches(2.0)
for a, b in ex:
    card(s, Inches(0.6), y, Inches(12.1), Inches(0.85))
    tf = tb(s, Inches(0.95), y, Inches(6.6), Inches(0.85), MSO_ANCHOR.MIDDLE)
    add_run(tf.paragraphs[0], a, 15, TXT)
    tf = tb(s, Inches(7.8), y, Inches(4.7), Inches(0.85), MSO_ANCHOR.MIDDLE)
    add_run(tf.paragraphs[0], '→  ' + b, 15, CORAL, bold=True)
    y += Inches(0.98)
footer(s, num(), 'Level 3 · 확장하기')

# ══════════════ 33. 실전 팁 ══════════════
s = new_slide()
title(s, '실전 팁 5가지')
bullets(s, [
    [('① 작게 시작해서 신뢰 쌓기', TXT, True), (' — 작은 수정 → 기능 → 리팩터링 순서로', DIM, False)],
    [('② 검증 방법을 함께 주기', TXT, True), (' — “고치고 테스트도 돌려서 확인해줘”', DIM, False)],
    [('③ 중간에 개입하기', TXT, True), (' — 이상하면 Esc, 방향만 다시 잡아주면 됩니다', DIM, False)],
    [('④ 커밋을 자주', TXT, True), (' — 잘 된 상태를 저장해두면 실험이 두렵지 않습니다', DIM, False)],
    [('⑤ 안 되는 이유를 묻기', TXT, True), (' — “왜 이 방법을 썼어?” 답을 듣고 배우세요', DIM, False)],
], Inches(0.7), Inches(2.1), Inches(12), size=18, gap=16)
footer(s, num(), '실전 팁')

# ══════════════ 34. 안티패턴 ══════════════
s = new_slide()
title(s, '이것만은 피하세요')
bullets(s, [
    [('diff 안 읽고 전부 승인', YELL, True), (' — 이해 못 한 코드가 쌓이면 결국 내 빚입니다', TXT, False)],
    [('한 대화에서 주제 계속 갈아타기', YELL, True), (' — 정확도가 떨어집니다. /clear 하세요', TXT, False)],
    [('모호한 불만 반복', YELL, True), (' — “아니 그거 말고” 대신 원하는 결과를 다시 설명', TXT, False)],
    [('비밀키·개인정보 붙여넣기', YELL, True), (' — .env 값은 가리고 공유하세요', TXT, False)],
    [('백업 없이 대규모 변경', YELL, True), (' — 커밋 먼저, 실험은 그 다음', TXT, False)],
], Inches(0.7), Inches(2.1), Inches(12), size=18, gap=16)
footer(s, num(), '실전 팁')

# ══════════════ 35. 치트시트 ══════════════
s = new_slide()
title(s, '치트시트 — 이것만 기억하세요')
data = [('claude', '시작'), ('claude -c', '마지막 대화 이어가기'),
        ('/init', 'CLAUDE.md 생성'), ('/clear · /compact', '초기화 · 요약'),
        ('Shift+Tab', '권한 모드 전환 (Plan Mode)'), ('Esc · Esc Esc', '중단 · 되돌리기'),
        ('/model · /status', '모델 변경 · 상태 확인'), ('claude mcp add', 'MCP 서버 연결')]
y = Inches(1.95)
for i in range(0, len(data), 2):
    for j in range(2):
        if i + j >= len(data): break
        k, v = data[i + j]
        x = Inches(0.6) + j * Inches(6.25)
        card(s, x, y, Inches(6.0), Inches(1.0))
        tf = tb(s, x + Inches(0.3), y, Inches(2.9), Inches(1.0), MSO_ANCHOR.MIDDLE)
        add_run(tf.paragraphs[0], k, 15, CORAL, bold=True, font=MONO)
        tf = tb(s, x + Inches(3.15), y, Inches(2.7), Inches(1.0), MSO_ANCHOR.MIDDLE)
        add_run(tf.paragraphs[0], v, 14, TXT)
    y += Inches(1.15)
footer(s, num(), '부록')

# ══════════════ 36. 리소스 ══════════════
s = new_slide()
title(s, '다음 단계 & 리소스')
bullets(s, [
    [('오늘', CORAL, True), ('  Level 1 실습 — 내 프로젝트에서 질문 3개, 수정 1건', TXT, False)],
    [('이번 주', CORAL, True), ('  /init + Plan Mode로 실제 업무 1건 완료', TXT, False)],
    [('이번 달', CORAL, True), ('  팀 CLAUDE.md 정착, 커스텀 커맨드 1개 만들기', TXT, False)],
    [('공식 문서  ', TXT, True), ('code.claude.com/docs', CORAL, False)],
    [('튜토리얼 영상  ', TXT, True), ('EP1~EP8 — 이 자료와 같은 커리큘럼으로 제작', CORAL, False)],
], Inches(0.7), Inches(2.2), Inches(12), size=18, gap=16)
footer(s, num(), '마무리')

# ══════════════ 37. Q&A ══════════════
s = new_slide(); num()
tf = tb(s, Inches(0.7), Inches(2.6), Inches(12), Inches(1.5))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
add_run(p, '질문 있으신가요?', 40, TXT, bold=True)
tf = tb(s, Inches(0.7), Inches(4.0), Inches(12), Inches(1.0))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
add_run(p, '가장 좋은 시작은 지금 터미널을 열고  ', 18, DIM)
add_run(p, 'claude', 18, CORAL, bold=True)
add_run(p, '  라고 입력하는 것입니다.', 18, DIM)

prs.save(OUT)
print(f'saved: {OUT} · {N} slides')
